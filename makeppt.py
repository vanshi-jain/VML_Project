from pptx import Presentation 
from pptx.util import Inches
import os
from os.path import isfile, join
from os import listdir
import glob
path = './smoothed-vit/'
raw_path = path+'raw/'
att_path = path+'attacked/'

prs = Presentation() 
blank_slide_layout = prs.slide_layouts[6] 
slide = prs.slides.add_slide(blank_slide_layout)

rawfiles = [f for f in listdir(raw_path) if isfile(join(raw_path, f))]
patchfiles = [f for f in listdir(att_path) if isfile(join(att_path, f))]
id_files = {k:[] for k in range(len(rawfiles))}
# print(id_files)
for i in range(len(rawfiles)):
    st = (rawfiles[i].split('_'))
    st2 = st[-1].split('.')
    id = int(st2[0])
    # print(id)
    if len(id_files[id]) > 0:
        t= id_files[id]
        id_files[id] = t + [raw_path+rawfiles[i]] 
    else:
        id_files[id] = [raw_path+rawfiles[i]] 

    st = (patchfiles[i].split('_'))
    st2 = st[-1].split('.')
    id2 = int(st2[0])
    if len(id_files[id2]) > 0:
        t = id_files[id2]
        id_files[id2] = t +[att_path+patchfiles[i]]
    else:
        id_files[id2] = [att_path+patchfiles[i]]

for i in range(0,len(rawfiles)-1,2):
    slide = prs.slides.add_slide(blank_slide_layout)
    pic = slide.shapes.add_picture(id_files[i][0], Inches(1.0), top=Inches(0.5),
                               width=Inches(3), height=Inches(3))
    pic = slide.shapes.add_picture(id_files[i][1], Inches(4.5), top=Inches(0.5),
                               width=Inches(3), height=Inches(3))
    
    pic = slide.shapes.add_picture(id_files[i+1][0], Inches(1.0), top=Inches(4.0),
                               width=Inches(3), height=Inches(3))
    pic = slide.shapes.add_picture(id_files[i+1][1], Inches(4.5), top=Inches(4.0),
                               width=Inches(3), height=Inches(3))
prs.save('test.pptx')
