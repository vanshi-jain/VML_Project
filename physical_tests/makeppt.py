# The script creates a PowerPoint presentation with images from two directories.
from pptx import Presentation
from pptx.util import Inches
import os
from os.path import isfile, join
from typing import Dict, List


def get_image_files(path: str) -> List[str]:
    """Returns list of file names in a directory."""
    return [f for f in os.listdir(path) if isfile(join(path, f))]


def extract_image_id(filename: str) -> int:
    """Extracts ID from a file name of the format 'xxx_<id>.ext'."""
    try:
        return int(filename.split('_')[-1].split('.')[0])
    except (IndexError, ValueError):
        raise ValueError(f"Filename format is incorrect: {filename}")


def organize_files_by_id(raw_files: List[str], patch_files: List[str], raw_path: str, att_path: str) -> Dict[int, List[str]]:
    """Organizes image file paths by their ID."""
    id_files = {}

    for file_list, base_path in [(raw_files, raw_path), (patch_files, att_path)]:
        for filename in file_list:
            img_id = extract_image_id(filename)
            full_path = join(base_path, filename)
            id_files.setdefault(img_id, []).append(full_path)

    return id_files


def add_images_to_slide(slide, image_paths: List[str], top_offset: float):
    """Adds a pair of images to the slide at a given vertical position."""
    slide.shapes.add_picture(image_paths[0], Inches(1.0), Inches(top_offset), width=Inches(3), height=Inches(3))
    slide.shapes.add_picture(image_paths[1], Inches(4.5), Inches(top_offset), width=Inches(3), height=Inches(3))


def create_presentation(id_files: Dict[int, List[str]], output_file: str):
    """Creates and saves a PowerPoint presentation with image pairs."""
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    sorted_ids = sorted(id_files.keys())
    for i in range(0, len(sorted_ids) - 1, 2):
        slide = prs.slides.add_slide(blank_slide_layout)
        add_images_to_slide(slide, id_files[sorted_ids[i]], top_offset=0.5)
        add_images_to_slide(slide, id_files[sorted_ids[i + 1]], top_offset=4.0)

    prs.save(output_file)
    print(f"Presentation saved as {output_file}")


def main():
    base_path = './smoothed-vit/'
    raw_path = join(base_path, 'raw/')
    att_path = join(base_path, 'attacked/')

    raw_files = get_image_files(raw_path)
    patch_files = get_image_files(att_path)

    id_files = organize_files_by_id(raw_files, patch_files, raw_path, att_path)
    create_presentation(id_files, 'test.pptx')


if __name__ == "__main__":
    main()
